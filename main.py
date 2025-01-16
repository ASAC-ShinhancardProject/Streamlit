import streamlit as st
# import anthropic
from langchain_core.messages import HumanMessage, AIMessage, SystemMessage
from langsmith import Client
from streamlit_feedback import streamlit_feedback
import os
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from pathlib import Path
import pandas as pd
from pptx import Presentation
import io
import time  # 추가된 부분
import docx
import base64
from pathlib import Path
from openai import OpenAI
from dotenv import load_dotenv
from PIL import Image

# 환경 변수 로드
load_dotenv()

# API 키 설정
openai_api_key = st.secrets.get("OPENAI_API_KEY") or os.getenv("OPENAI_API_KEY")

# OpenAI 설정
client = OpenAI(api_key=openai_api_key)

# Streamlit 앱 설정
st.set_page_config(page_title="집코기", page_icon=Image.open("챗봇.png"), layout="wide")

# Google Fonts에서 원하는 폰트 로드
st.markdown("""
<head>

<style>
[data-testid="stDecoration"] {
        background: none !important;
        height: 0px;
        margin: 0 !important;
        padding: 0 !important;}

@import url('https://fonts.googleapis.com/css2?family=Pretendard:wght@300;400;500;700&display=swap');

/* (1) 상단 헤더 스타일 */
    .chatbot-header {
        /*background-color: #4f46e5 !important;*/
        background-color: #29caff !important;
        
        color: #fff !important;
        text-align: center !important;
        padding: 15px 20px !important;
        font-size: 1.2rem !important;
        font-weight: bold !important;
        border-radius: 12px 12px 0 0 !important;
        margin: 0 auto !important;
        box-shadow: 0 4px 10px rgba(0,0,0,0.1) !important;
        width: 100% !important;
        font-family: 'Pretendard', sans-serif !important;
    }

    /* (2) st.chat_message + st.chat_input 전체 컨테이너 */
    [data-testid="stChatContainer"] {
        margin: 0 auto !important;
        width: 100% !important;
        background-color: #fff !important;
        border-left: 1px solid #ddd !important;
        border-right: 1px solid #ddd !important;
        border-bottom: 1px solid #ddd !important;
        border-radius: 0 0 12px 12px !important;
        box-shadow: 0 4px 10px rgba(0,0,0,0.1) !important;
        display: flex !important;
        flex-direction: column !important;
        justify-content: space-between !important;
        overflow: hidden !important;
        position: relative !important;
        font-family: 'Pretendard', sans-serif !important;
    }

    /* (3) 메시지 영역 */
    [data-testid="stChatMessages"] {
        flex: 1 !important;
        padding: 20px !important;
        overflow-y: auto !important;
        background-color: #fff !important;
        font-family: 'Pretendard', sans-serif !important;
    }

    /* (4) 입력창 영역 */
    [data-testid="stChatInput"] {
        margin: 0 20px !important;
        border: 1px solid #ddd !important;
        border-radius: 0 0 12px 12px !important;
        padding: 10px !important;
        background-color: #f9f9f9 !important;
        display: flex !important;
        align-items: center !important;
        position: relative !important;
        box-sizing: border-box !important;
        overflow: hidden !important;
        font-family: 'Pretendard', sans-serif !important;
    }
    /* 화살표 버튼 위치 조정 */
    [data-testid="stChatInput"] button {
        position: absolute !important;
        right: 15px !important;
        top: 50% !important;
        transform: translateY(-50%) !important;
        background-color: transparent !important;
        border: none !important;
        box-shadow: none !important;
        outline: none !important;
        padding: 0 !important;
        margin: 0 !important;
        cursor: pointer !important;
    }
    /* 화살표 아이콘 색상 */
    [data-testid="stChatInput"] button svg {
        fill: #666 !important;
    }

    /* (5) 사용자/AI 메시지 (배경색/테두리) */
    [data-testid="stChatMessage"][data-testid="stChatMessageRole=assistant"] {
        background-color: #f7f9ff !important; /* 은은한 파란톤 */
        border: 1px solid #ececec !important;
        margin-bottom: 1rem !important;
        padding: 1rem !important;
        border-radius: 8px !important;
        font-family: 'Pretendard', sans-serif !important;
    }
    [data-testid="stChatMessage"][data-testid="stChatMessageRole=user"] {
        background-color: #ffffff !important;
        border: 1px solid #ececec !important;
        margin-bottom: 1rem !important;
        padding: 1rem !important;
        border-radius: 8px !important;
        font-family: 'Pretendard', sans-serif !important;
    }
    </style>
    """,
    unsafe_allow_html=True
)


# 현재 스크립트의 디렉토리를 기준으로 assets 폴더 경로 설정
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
ASSETS_DIR = os.path.join(SCRIPT_DIR, 'static')

def get_avatar_path(role: str) -> str:
    """이미지 파일의 절대 경로를 반환"""
    image_path = os.path.join(ASSETS_DIR, f'{role}_character.png')
    if os.path.exists(image_path):
        return image_path
    print(f"Warning: Image not found at {image_path}")  # 디버깅용
    return None

def send_message(message, role, save=True):
    """Display message with appropriate avatar"""
    avatar_path = get_avatar_path('human' if role == 'human' else 'bot')
    try:
        with st.chat_message(role, avatar=avatar_path):
            st.markdown(message, unsafe_allow_html=True)
        if save:  # 메시지를 한 번만 저장
            save_message(message, role)
    except Exception as e:
        print(f"Error displaying message with avatar: {e}")
        with st.chat_message(role):
            st.markdown(message, unsafe_allow_html=True)
        if save:
            save_message(message, role)

def get_image_as_base64(image_path):
    """이미지를 Base64 문자열로 변환"""
    try:
        with open(image_path, "rb") as img_file:
            return base64.b64encode(img_file.read()).decode("utf-8")
    except Exception as e:
        st.error(f"이미지를 불러오는 중 오류가 발생했습니다: {e}")
        return ""
    

# 배경 이미지 추가
bg_image_path = "static/bg.png"  # 배경 이미지 경로

if Path(bg_image_path).exists():
    bg_image_base64 = get_image_as_base64(bg_image_path)
    st.markdown(
        f"""
        <style>
        /* 전체 페이지 배경 */
        html {{
            background-image: url("data:image/png;base64,{bg_image_base64}");
            background-size: cover;
            background-attachment: fixed;
            background-repeat: no-repeat;
        }}

        /* 텍스트 입력창 하단 영역 (stChatInput) */
        [data-testid="stApp"]{{
            background-image: url("data:image/png;base64,{bg_image_base64}");
            background-size: cover;
    /       background: rgba(255, 255, 255, 0); /* 투명화 *
        }}
        
        /* 사이드바 배경 투명화 */
        [data-testid="stSidebar"] {{
            background: rgba(255, 255, 255, 0); /* 투명화 */
        }}
        
        /* 사이드바 배경 투명화 */
        [data-testid="stHeader"] {{
            background: rgba(255, 255, 255, 0); /* 투명화 */
        }}
        
        /* 사이드바 배경 투명화 */
        [data-testid="stBottom"] {{
            background: rgba(255, 255, 255, 0); /* 투명화 */
        }}
        
        /* 사이드바 배경 투명화 */
        [data-testid="stBottom"] > div {{
            background: rgba(255, 255, 255, 0); /* 투명화 */
        }}
        
        /*특정 영역 색상 변경 */
        .stMain {{
            background: rgba(255, 255, 255, 255); /* 투명화 */
        }}
        """,
        unsafe_allow_html=True
    )
else:
    st.warning("배경 이미지 파일이 존재하지 않습니다.")



# 페이지 제목
# st.markdown('<h1 class="custom-title">🏡 집코기</h1>', unsafe_allow_html=True)
#st.markdown('<h3 class="custom-title1"> 아파트에 관해 물어보세요 (아파트 정보, 청약, 대출, 세금 등) </h3>', unsafe_allow_html=True)

st.markdown('<div class="chatbot-header">집코기</div>', unsafe_allow_html=True)


if not openai_api_key:
    st.error("OpenAI API key 설정을 확인해주시기 바랍니다.")
    st.stop()

chatbot_content = """
1. **[신혼부부 프롬프트]**의 **Q1**해당된다면 **대출 상품**, **주체**, **대출 조건**, **특징** ,**금리**, **대출한도**, **링크**, **LTV**, **DTI**에 대해서 요약없이 그대로 얘기해줍니다.

**[신혼부부 프롬프트]**
###**Q1:** 신혼부부인데, 대출 상품 알려줘.

### [대출 상품 1: 디딤돌 대출]
- **주체:** 한국주택금융공사, 주택도시기금 기금e든든 홈페이지 또는 기금수탁은행(국민, 농협, 신한, 우리, 하나, 부산, 대구)으로 일원화되어 운영
- **대출 조건:** 
  - 세대원 전원 무주택 세대주
  - 부부합산 연소득  6,000만 원 이하(신혼가구 8,500만 원, 생애최초 주택 구입자•다자녀•2자녀가구는 7,000만 원)
  - 85㎡ 이하 주택 담보주택 평가액 5억 원 이하(신혼가구 및 2자녀 이상 가구 6억 원)
  - 주택가격(매매가) 기준 제한 존재 (수도권 5억 원, 비수도권 3억 원 등)
  - 본인 및 배우자 합산 순자산 가액 4.88억 원 이하
- **특징:** 일반 시중은행 대출보다 금리 우대 혜택이 있음.
- **금리:** 대체로 2.65%~3.95% (고정금리 또는 5년 단위 변동금리)
- **대출한도:** 2.5억 원 이내(생애최초 주택구입자 3억 원, 신혼부부•다자녀•2자녀가구는 4억 원) 
- **링크:** https://www.hf.go.kr/ko/sub01/sub01_02_01.do#!#close

### [대출 상품 2: 보금자리론]
- **주체:** 한국주택금융공사
- **대출 조건:** 
  - 무주택 OR 1주택
  - 부부 합산 연소득 7,000만 원 이하(신혼가구 8,500만 원, 1자녀가구 8,000만 원, 2자녀가구 9,000만 원, 다자녀가구 1억 원, 전세사기피해자 등 없음)
  - 주택 가격 6억 이내
- **특징:** 고정금리 형태, 만기 10~50년 등 선택 가능
- **금리:** 시기에 따라 변동되지만, 시중은행 대비 다소 낮거나 비슷한 수준의 고정금리를 제공
- **대출한도:** 3.6억 원 이내(다자녀가구•전세사기피해자등 4억 원, 생애최초 주택구입자 4.2억 원)
- **LTV:**
  - 비조정대상지역: 아파트는 70%, 기타주택(연립•다세대•단독주택)은 65%
  - 조정대상지역: 아파트는는 60%, 기타주택(연립•다세대•단독주택)은 55% 
- **DTI:**
  - 비조정대상지역: 60% 이내
  - 조정대상지역: 10%P 차감하여 적용
- **링크:** https://www.hf.go.kr/ko/sub01/sub01_01_01.do

### [대출 상품 3: 신혼부부전용 구입자금]
- **주체:** 주택도시기금, 지자체, 시중은행 협약상품 등
- **대출 조건:** 
  - 혼인신고일로부터 7년 이내의 신혼부부 또는 3개월 이내 결혼예정자
  - 무주택 부부(한쪽이라도 주택 보유 시 혜택 제한)
  - 부부 합산 연소득 8,500만 원 이하
  - 본인 및 배우자 합산 순자산 가액 4.88억 원 이하
- **특징:** 금리 우대(최대 0.2~0.3%p 정도 추가 할인), LTV/DTI 완화 등 혜택
- **금리:** 연 2.35% ∼ 연 3.65% 
- **링크:** https://nhuf.molit.go.kr/FP/FP05/FP0503/FP05030601.jsp
"""

chatbot_content = """
1. **[프롬프트]**의 **Q1**해당된다면 **사용자 유형**, **본인의 연령대** 정보를 요약하지 말고 모든 세부 정보를 포함하여 답변하세요.
2. **[신혼부부 프롬프트]**의 **Q1**해당된다면 **[신혼부부 프롬프트 Q1의 답변]**에 있는 **대출 상품**, **주체**, **대출 조건**, **특징** ,**금리**, **대출한도**, **링크**, **LTV**, **DTI**에 대해서 요약하지 말고 모든 세부 정보를 포함하여 답변하세요. 
3. **[신혼부부 프롬프트]**의 **Q1**해당된다면 **[신혼부부 프롬프트 Q1의 답변]**에 있는 **아파트**, **주소**, **세대수**, **저/최고층**, **총주차대수**, **용적률**, **건폐율**, **건설사**를 요약하지 말고 모든 세부 정보를 포함하여 답변하세요. 
4. **[신혼부부 프롬프트]**의 **Q2**해당된다면 **[신혼부부 프롬프트 Q2의 답변]**에 있는 **아파트**, **주소**, **교통**, **연식**, **학군**, **난방 방식**, **종합 분석**, **링크**에 대해서 요약없이 그대로 얘기해줍니다. 
5. **[중장년층 프롬프트]**의 **Q1**해당된다면 **[중장년층 프롬프트 Q1의 답변]**에 있는 **대출 상품**, **주체**, **대출 조건**, **특징** ,**금리**, **대출한도**, **링크**, **LTV**, **DTI**에 대해서 요약하지 말고 모든 세부 정보를 포함하여 답변하세요. 
6. **[중장년층 프롬프트]**의 **Q1**해당된다면 **[중장년층 프롬프트 Q1의 답변]**에 있는 **아파트**, **주소**, **세대수**, **저/최고층**, **총주차대수**, **용적률**, **건폐율**, **건설사**를 요약하지 말고 모든 세부 정보를 포함하여 답변하세요. 
4. **[중장년층 프롬프트]**의 **Q2**해당된다면 **[중장년층 프롬프트 Q2의 답변]**에 있는 **아파트**, **주소**, **교통**, **연식**, **학군**, **난방 방식**, **종합 분석**, **링크**에 대해서 요약없이 그대로 얘기해줍니다. 

**[프롬프트]**
**Q1:** 아파트를 추천해주세요.

### [사용자 유형]
부동산 금융 서비스를 추천하기 위해 몇 가지 질문을 드릴게요. 아래 내용을 알려주세요:
1. **본인의 연령대**: 청년층, 신혼부부, 중장년층, 시니어층 중 어디에 해당하시나요?
2. **무주택 여부**: 현재 무주택자이신가요?
3. **부부 합산 소득**: 대략적인 연 소득을 알려주세요.
4. **추천받고 싶은 지역**: 어떤 지역에서 아파트를 추천받고 싶으신가요?

이 정보를 기반으로 적합한 금융 상품과 주택 추천 정보를 제공해 드리겠습니다!

**[신혼부부 프롬프트]**
**Q1:** 저는 신혼부부로, 무주택자이며 부부 합산 소득이 6,000만 원 이하입니다. 관악, 광명, 안양 지역을 중심으로 아파트를 추천받고 싶습니다.
**Q2:** 추천받은 대출 상품은 만족스럽습니다. 안내받은 아파트 중에서 '철산광복현대', '더포레스트힐아파트', '영등포중흥에스클래스'에 대해 더 자세한 정보를 알고 싶습니다. 

**[중장년층 프롬프트]**
**Q1:** 저는 중장장년층으로,  1주택자이며 부부 합산 소득이 7,000만 원 이하입니다. 영등포구, 구로구, 광명시 지역을 중심으로 아파트를 추천받고 싶습니다.
**Q2:** 추천받은 대출 상품은 만족스럽습니다. 안내받은 아파트 중에서 '구로두산위브', '광명아크포레자이위브', '다산해모로아파트'에 대해 더 자세한 정보를 알고 싶습니다.


**[신혼부부 프롬프트 Q1의 답변]**
신혼부부로서 무주택 상태이며, 부부 합산 소득이 6,000만 원 이하이신 상황에서 대출 상품과 관악구, 광명시, 안양시 지역의 아파트를 안내해드리겠습니다. 

먼저, 조건에 맞는 대출 상품과 아파트 옵션은 다음과 같습니다. 
### [대출 상품 1: 디딤돌 대출]
- **주체:** 한국주택금융공사, 주택도시기금 기금e든든 홈페이지 또는 기금수탁은행(국민, 농협, 신한, 우리, 하나, 부산, 대구)으로 일원화되어 운영
- **대출 조건:** 
  - 세대원 전원 무주택 세대주
  - 부부합산 연소득  6,000만 원 이하(신혼가구 8,500만 원, 생애최초 주택 구입자•다자녀•2자녀가구는 7,000만 원)
  - 85㎡ 이하 주택, 담보주택 평가액 5억 원 이하(신혼가구 및 2자녀 이상 가구 6억 원)
  - 주택가격(매매가) 기준 제한 존재 (수도권 5억 원, 비수도권 3억 원 등)
  - 본인 및 배우자 합산 순자산 가액 4.88억 원 이하
- **특징:** 일반 시중은행 대출보다 금리 우대 혜택 
- **금리:** 대체로 2.65%~3.95% (고정금리 또는 5년 단위 변동금리)
- **대출한도:** 2.5억 원 이내(생애최초 주택구입자 3억 원, 신혼부부•다자녀•2자녀가구는 4억 원) 
- 자세한 내용은 **링크**를 통해 확인하시기 바랍니다. https://www.hf.go.kr/ko/sub01/sub01_02_01.do#!#close

### [대출 상품 2: 보금자리론]
- **주체:** 한국주택금융공사
- **대출 조건:** 
  - 무주택 OR 1주택
  - 부부 합산 연소득 7,000만 원 이하(신혼가구 8,500만 원, 1자녀가구 8,000만 원, 2자녀가구 9,000만 원, 다자녀가구 1억 원, 전세사기피해자 등 없음)
  - 주택 가격 6억 이내
- **특징:** 고정금리 형태, 만기 10~50년 등 선택 가능
- **금리:** 시기에 따라 변동되지만, 시중은행 대비 다소 낮거나 비슷한 수준의 고정금리를 제공
- **대출한도:** 3.6억 원 이내(다자녀가구•전세사기피해자등 4억 원, 생애최초 주택구입자 4.2억 원)
- **LTV:**
  - 비조정대상지역: 아파트는 70%, 기타주택(연립•다세대•단독주택)은 65%
  - 조정대상지역: 아파트는는 60%, 기타주택(연립•다세대•단독주택)은 55% 
- **DTI:**
  - 비조정대상지역: 60% 이내
  - 조정대상지역: 10%P 차감하여 적용
- 자세한 내용은 **링크**를 통해 확인하시기 바랍니다. https://www.hf.go.kr/ko/sub01/sub01_01_01.do

### [대출 상품 3: 신혼부부전용 구입자금]
- **주체:** 주택도시기금, 지자체, 시중은행 협약상품 등
- **대출 조건:** 
  - 혼인신고일로부터 7년 이내의 신혼부부 또는 3개월 이내 결혼예정자
  - 무주택 부부(한쪽이라도 주택 보유 시 혜택 제한)
  - 부부 합산 연소득 8,500만 원 이하
  - 본인 및 배우자 합산 순자산 가액 4.88억 원 이하
- **특징:** 금리 우대(최대 0.2~0.3%p 정도 추가 할인), LTV/DTI 완화 등 혜택
- **금리:** 연 2.35% ∼ 연 3.65% 
- 자세한 내용은 **링크**를 통해 확인하시기 바랍니다. https://nhuf.molit.go.kr/FP/FP05/FP0503/FP05030601.jsp

---

### [아파트 1: 신림주공1단지]
- **주소:** 서울특별시 관악구 신림동 1704-1 
- **세대수:** 90세대(총1개동)
- **저/최고층:** 21층/23층
- **총주차대수:** 353대
- **용적률:** 255%
- **건폐율:** 16%
- **건설사:** 대한주택공사

### [아파트 2: 철산광복현대]
- **주소:** 경기도 광명시 철산동 55-1
- **세대수:** 841세대(총7개동)
- **저/최고층:** 9층/23층
- **총주차대수:** 750대
- **용적률:** 317%
- **건폐율:** 22%
- **건설사:** 현대건설(주)

### [아파트 3: 더포레스트힐아파트]
- **주소:** 경기도 안양시 동안구 비산동 1155
- **세대수:** 2044세대(총18개동)
- **저/최고층:** 12층/25층
- **총주차대수:** 2061대
- **용적률:** 225%
- **건폐율:** 16%
- **건설사:** (주)국제종합토건

---

다음은 조건에 부합하지 않지만, 함께 고려해볼 만한 아파트입니다. 
### [아파트 1: 영등포중흥에스클래스]
- **주소:** 서울특별시 영등포구 양평동1가 271 
- **세대수:** 308세대(총2개동)
- **저/최고층:** 24층/24층
- **총주차대수:** 374대
- **용적률:** 399%
- **건폐율:** 46%
- **건설사:** 중흥토건(주)

### [아파트 2: 세운푸르지오헤리시티]
- **주소:** 서울특별시 중구 인현동2가 240
- **세대수:** 321세대(기타임대 40세대 포함, 총1개동)
- **저/최고층:** 15층/26층
- **총주차대수:** 182대
- **용적률:** 937%
- **건폐율:** 60%
- **건설사:** (주)대우건설

### [아파트 3: 래미안안양메가트리아]
- **주소:** 경기도 안양시 만안구 안양동 1393
- **세대수:** 4250세대(총35개동)
- **저/최고층:** 17층/32층
- **총주차대수:** 4933대
- **용적률:** 244%
- **건폐율:** 16%
- **건설사:** 삼성물산(주)

위의 아파트와 대출 상품을 고려하여 주거 계획을 세우시면 도움이 될 것입니다. 

**[중장년층 프롬프트 Q1의 답변]**
중장년층로서 1주택자이며, 부부 합산 소득이 7,000만 원 이하이신 상황에서 대출 상품과 관악구, 광명시, 안양시 지역의 아파트를 안내해드리겠습니다. 

먼저, 조건에 맞는 대출 상품과 아파트 옵션은 다음과 같습니다. 
### [대출 상품 1: 주택연금 사전예약 보금자리론]
- **주체:** 한국주택금융공사
- **대출 조건:** 
  - 본인 또는 배우자가 만 40세 이상인 경우 신청 가능
  - 대출신청 시 본건 담보주택 외 1주택을 기준으로 추가 주택을 보유할 경우 처분기한 내 처분
  - 부부합산 연소득  7,000만 원 이하
  - 담보주택 평가액 5억 원 이하
- **특징:** 만 55세 이후에는 대출 잔여만기를 불문하고 주택연금으로 전환 시 전환장려금 형태로 지급
- **금리:** 
  - 대출 전액상환일까지 기본 우대금리 0.15%p가 적용
  - 기존대출(‘15.12.31일 이전 실행건)이 변동금리 또는 일시상환 대출인 경우에는 0.15%p가 추가되어 총 0.3%p의 우대금리가 적용
- **대출한도:**  2억원 이하 대출실행금액에 대해서만 전환장려금이 지급
- 자세한 내용은 **링크**를 통해 확인하시기 바랍니다. https://www.hf.go.kr/ko/sub01/sub01_01_02.do#!#m7

---

### [아파트 1: 영등포중흥에스클래스]
- **주소:** 서울특별시 영등포구 양평동1가 271 
- **세대수:** 308세대(총2개동)
- **저/최고층:** 24층/24층
- **총주차대수:** 374대
- **용적률:** 399%
- **건폐율:** 46%
- **건설사:** 중흥토건(주)


### [아파트 2: 구로두산위브]
- **주소:** 서울시 구로구 구로동 1278
- **세대수:** 660세대(총9개동)
- **저/최고층:** 9층/21층
- **총주차대수:** 579대
- **용적률:** 252%
- **건폐율:** 18%
- **건설사:** 두산산업개발(주)

### [아파트 3: 광명아크포레자이위브]
- **주소:** 경기도 광명시 광명동 787-4
- **세대수:** 2104세대(기타임대 101세대 포함, 총19개동)
- **저/최고층:** 14층/29층
- **총주차대수:** 2431대
- **용적률:** 279%
- **건폐율:** 17%
- **건설사:** 지에스건설(주)

---

다음은 조건에 부합하지 않지만, 함께 고려해볼 만한 아파트입니다. 

### [아파트 1: 평촌자이아이파크]
- **주소:** 경기도 안양시 동안구 비산동 1185
- **세대수:** 2737세대(기타임대 146세대 포함, 총22개동)
- **저/최고층:** 18층/29층
- **총주차대수:** 3626대
- **용적률:** 248%
- **건폐율:** 17%
- **건설사:** 지에스건설(주)

### [아파트 2: 다산해모로]
- **주소:** 경기도 남양주시 다산동 6297
- **세대수:** 449세대(기타임대 37세대 포함, 총5개동)
- **저/최고층:** 16층/29층
- **총주차대수:** 517대
- **용적률:** 278%
- **건폐율:** 17%
- **건설사:** (주)한진중공업

### [아파트 3: 정릉꿈에그린]
- **주소:** 서울시 성북구 정릉동 1037
- **세대수:** 349세대(기타임대 60세대 포함, 총8개동)
- **저/최고층:** 2층/19층
- **총주차대수:** 395대
- **용적률:** 257%
- **건폐율:** 33%
- **건설사:** (주)한화건설

**[신혼부부 프롬프트 Q2의 답변]**

문의주신 "철산광복현대", "더포레스트힐아파트", "영등포중흥에스클래스"에 대한 자세한 정보를 안내드리겠습니다. 

### [아파트 1: 철산광복현대]
- **주소:** 경기도 광명시 철산동 55-1
- **교통:** 
  - 7호선 철산역 및 광명사거리역 인근 (1km 내 총 2개 역 접근 가능)
  - 2호선, 7호선, 경부선, 경인선 등 (3km 내)
- **연식:** 준공 1998년, 약 25년차, 중고 아파트
- **학군:**
  - 유치원: 2개(아파트 단지내)
  - 초등학교: 우수, 도보 10분 이하 접근 가능
  - 중학교: 매우 우수, 도보 10분 이하 접근 가능
- **난방 방식:** 중앙난방 
- **종합 분석:** 
  - 교통 접근성이 양호하며, 7호선과 주요 노선 접근성이 좋음.
  - 학군 환경이 우수하며, 초등학교와 중학교 모두 도보 거리에 있어 학부모에게 매력적.
  - 1차 의료기관 밀도가 적당히 높아 생활 편의성이 양호함.
  - 상위 10위 건설사 '현대건설'의 브랜드에 해당함.
- 자세한 내용은 **링크**를 통해 확인하시기 바랍니다. https://public.tableau.com/app/profile/dyjeong/viz/APT_EDA/Title?publish=yes&P_SGG=광명시&P_APT_NM=철산광복현대

### [아파트 2: 더포레스트힐아파트]
- **주소:** 경기도 안양시 동안구 비산동 1155
- **교통:** 
  - 경부선 관악역, 명학역, 안양역 및 과천선 범계역 인근 (3km 내)
- **연식:** 준공 2003년, 약 22년차, 중고 아파트
- **학군:**
  - 유치원: 1개(아파트 단지내)
  - 초등학교: 매우 미흡, 도보 30분 초과
  - 중학교: 매우 미흡, 도보 30분 초과
- **난방 방식:** 개별난방 
- **종합 분석:** 
  - 경부선과 과천선 접근성이 양호하여 대중교통 이용이 편리함.
  - 학군 및 교육 환경은 다소 부족하며, 초·중학교 접근성 개선이 필요함.
  - 1차 의료기관 밀도가 적당히 높아 생활 편의성이 양호함.
  - 개별난방 방식으로 선호도는 다소 낮으며, 중고 아파트로 실거주보다는 투자용으로 적합할 가능성 존재함.
- 자세한 내용은 **링크**를 통해 확인하시기 바랍니다. https://public.tableau.com/app/profile/dyjeong/viz/APT_EDA/Title?publish=yes&P_SGG=동안구&P_APT_NM=더포레스트힐아파트

### [아파트 3: 영등포중흥에스클래스]
- **주소:** 서울특별시 영등포구 양평동1가 271
- **교통:** 
  - 5호선 양평역 (500m, 도보 7분)
  - 2호선, 5호선, 경부선, 9호선 등 (3km 내)
- **연식:** 준공 2021년, 약 3년차, 준신축 아파트
- **학군:**
  - 초등학교: 우수, 도보 10분 이하 접근 가능
  - 중학교: 매우 우수, 도보 10분 이하 접근 가능
- **난방 방식:** 지역난방 
- **종합 분석:** 
  - 교통 접근성이 양호하며, 5호선 양평역과 인접해 출퇴근이 편리함.
  - 초·중학교와 가까운 거리에 위치해 학군 환경이 우수하며, 교육적 장점이 큼.
  - 1차 의료기관 밀도가 높아 생활 편의성이 우수.
  - 준신축 아파트로, 재건축 가능성은 낮으나 난방 방식과 시설이 최신.
- 자세한 내용은 **링크**를 통해 확인하시기 바랍니다. https://public.tableau.com/app/profile/dyjeong/viz/APT_EDA/Title?publish=yes&P_SGG=영등포구&P_APT_NM=영등포중흥에스클래스

이 아파트들은 각기 다른 위치와 특성을 가지고 있어, 해당 지역의 주거 환경이나 편의성에 따라 선택이 달라질 수 있습니다. 추가 요청 사항이 있다면 말씀해 주세요!

**[중장년층 프롬프트 Q2의 답변]**

문의주신 "구로두산위브", "광명아크포레자이위브", "다산해모로아파트"에 대한 자세한 정보를 안내드리겠습니다. 

### [아파트 1: 구로두산위브]
- **주소:** 서울시 구로구 구로동 1278
- **교통:** 
  - 7호선 남구로역 (500m, 도보 7분)
  - 2호선 대림역 및 신도림역, 경부선 구로역 등 다중 노선 접근 가능 (3km 내)
- **연식:** 준공 2006년, 약 19년차, 중고 아파트
- **학군:**
  - 초등학교: 매우 우수, 도보 5분 이하 접근 가능
  - 중학교: 매우 우수, 도보 10분 이하 접근 가능
- **난방 방식:** 개별난방 
- **종합 분석:** 
  - 교통 접근성이 매우 우수하며, 7호선 남구로역과 2호선 대림역 등 다중 노선 이용 가능.
  - 학군 및 교육 환경이 우수하며, 초·중학교 모두 도보 거리에 위치.
  - 1차 의료기관 밀도가 높아 생활 편의성이 뛰어나며, 발달된 상권으로 상업시설 접근성 우수.
- 자세한 내용은 **링크**를 통해 확인하시기 바랍니다. https://public.tableau.com/app/profile/dyjeong/viz/APT_EDA/Title?publish=yes&P_SGG=구로구&P_APT_NM=구로두산위브

### [아파트 2: 광명아크포레자이위브]
- **주소:** 경기도 광명시 광명동 787-4
- **교통:** 
  - 7호선 광명사거리역 (1km, 도보 14분)
  - 7호선, 경부선, 경인선 등 주요 노선 접근 가능 (3km 내)
- **연식:** 준공 2021년, 약 3년차, 준신축 아파트
- **학군:**
  - 초등학교: 매우 미흡, 도보 30분 초과
  - 중학교: 매우 미흡, 도보 30분 초과
- **난방 방식:** 지역난방
- **종합 분석:** 
  - 교통 접근성이 양호하며, 7호선과 경부선, 경인선 등 주요 노선과의 연결성이 좋음.
  - 1차 의료기관 밀도가 매우 높아 생활 편의성이 뛰어남.
  - 상위 10위 건설사 'GS건설'의 브랜드이고, 준신축 아파트로 시설과 난방 시스템에서 우수성이 기대됨.
- 자세한 내용은 **링크**를 통해 확인하시기 바랍니다. https://public.tableau.com/app/profile/dyjeong/viz/APT_EDA/Title?publish=yes&P_SGG=광명시&P_APT_NM=광명아크포레자이위브

### [아파트 3: 다산해모로아파트]
- **주소:** 경기도 남양주시 다산동 6297
- **교통:** 
  - 중앙선 구리역 (500m, 도보 7분)
  - 도농역 이용 가능(3km 내 총 유동인구: 1,406,025명)
- **연식:** 준공 2021년, 약 3년차, 준신축 아파트
- **학군:**
  - 초등학교: 우수, 도보 10분 이하 접근 가능
  - 중학교: 우수, 도보 10분 이하 접근 가능
- **난방 방식:** 지역난방
- **종합 분석:** 
  - 중앙선 구리역과 도농역과 가까워 교통 접근성이 우수.
  - 학군 환경이 우수하며 초등학교 및 중학교 모두 가까운 거리에 위치.
  - 발달 상권의 접근성이 뛰어나며, 준신축 아파트로서 최신 설비와 지역난방의 장점이 큼.
- 자세한 내용은 **링크**를 통해 확인하시기 바랍니다. https://public.tableau.com/app/profile/dyjeong/viz/APT_EDA/Title?publish=yes&P_SGG=남양주시&P_APT_NM=다산해모로아파트

이 아파트들은 각기 다른 위치와 특성을 가지고 있어, 해당 지역의 주거 환경이나 편의성에 따라 선택이 달라질 수 있습니다. 추가 요청 사항이 있다면 말씀해 주세요!
"""


# 세션 상태 초기화
if 'initialized' not in st.session_state:
    st.session_state.file_qa_messages = [SystemMessage(content = chatbot_content)]
    st.session_state.file_qa_content = None
    st.session_state.file_qa_data = None
    st.session_state.initialized = True
    st.session_state.greeting_displayed = False  # 인사 표시 여부를 추적하는 새로운 변수


# 사이드바에 파일 업로드 위젯 추가
with st.sidebar:
    # st.markdown("**[만든사람]** 황서연 이규희 장이철 정도영")
    #st.markdown("File Upload (Optional)")
    uploaded_file = st.file_uploader("Upload a file", type=("docx", "pdf", "csv", "xlsx", "pptx"))
    if uploaded_file:
        st.success(f"File '{uploaded_file.name}' uploaded successfully.")

# 채팅 초기화 버튼
reset_history = st.sidebar.button("채팅 초기화")

if reset_history:
    st.session_state.file_qa_messages = []
    st.session_state.file_qa_content = None
    st.session_state.greeting_displayed = False  # 초기화 시 인사 표시 상태도 초기화

# 파일 처리 함수
def process_file(file):
    if file is None:
        return None, None
    try:
        if file.type == "text/plain":
            content = file.read().decode()
            return content, None
        elif file.type == "application/pdf":
            pdf_reader = PdfReader(io.BytesIO(file.read()))
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text, None
        elif file.type == "text/csv":
            df = pd.read_csv(file)
            return df.to_string(), df
        elif file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            df = pd.read_excel(file)
            return df.to_string(), df
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = Presentation(file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + "\n"
            return text, None
        elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":  # 추가된 부분
            doc = docx.Document(file)
            text = ""
            for para in doc.paragraphs:
                text += para.text + "\n"
            return text, None
        else:
            return f"Unsupported file type: {file.type}", None
    except Exception as e:
        st.error(f"Error processing file: {str(e)}")
        return None, None

# 파일 업로드 및 처리
if uploaded_file:
    with st.spinner("Processing uploaded file..."):
        st.session_state.file_qa_content, st.session_state.file_qa_data = process_file(uploaded_file)
        if st.session_state.file_qa_content:
            st.sidebar.success(f"File '{uploaded_file.name}' processed successfully.")
        else:
            st.sidebar.error("Failed to process the file.")

# 챗봇 응답 함수
def get_chatbot_response(prompt, file_content=None):
    try:
        messages = [
            {"role": "system", "content": chatbot_content},
            {"role": "user", "content": prompt}
        ]

        response = client.chat.completions.create(
            model="gpt-4o",
            messages=messages,
            max_tokens=3200,
        )

        response_text = response.choices[0].message.content
        return response_text
    except Exception as e:
        st.error(f"Error in chatbot response: {str(e)}")
        return None

# 스트리밍 응답 함수
def stream_response(response_text, response_container):
    response = ""
    for char in response_text:
        response += char
        response_container.markdown(response, unsafe_allow_html=True)
        time.sleep(0.005)

# 사이드바에서 텍스트 분석 챗봇 모드 추가
# def create_sidebar_with_text_analysis():
#     """사이드바에서 텍스트 분석 챗봇 모드를 추가합니다."""
#     with st.sidebar:
#         st.markdown("### 🤖 챗봇 모드 선택")
        
#         # 모드 선택
#         mode = st.radio(
#             "원하시는 모드를 선택하세요:",
#             ["부동산 인사이트", "부동산 정책/금융"],
#             index=0,
#             key="chat_mode_sidebar"  # 고유 키로 변경
#         )
        
#         # 세션 상태 업데이트
#         st.session_state.analysis_mode = (mode == "텍스트 분석 챗봇")

# 인사 메시지 표시 (챗봇이 먼저 인사)
if not st.session_state.greeting_displayed:
    greeting = AIMessage(content=(
        "안녕하세요. 아파트 시장에 대해 물어보세요 \n"
        "내 집 마련과 시장 정보에 대한 질문 및 답변이 가능합니다."
    ))
    st.session_state.file_qa_messages.append(greeting)
    st.session_state.greeting_displayed = True


# 채팅 기록에서 가장 최근 메시지만 표시
# if st.session_state.file_qa_messages:
#     last_message = st.session_state.file_qa_messages[-1]
#     if isinstance(last_message, AIMessage):
#         st.chat_message("ai", avatar="챗봇.png").write(last_message.content)
#     elif isinstance(last_message, HumanMessage):
#         st.chat_message("human", avatar="질문.png").write(last_message.content)

for msg in st.session_state.file_qa_messages:
    if isinstance(msg, SystemMessage):
        continue  # 시스템 메시지는 숨김
    elif isinstance(msg, AIMessage):
        st.chat_message("ai", avatar="챗봇.png").write(msg.content)
    elif isinstance(msg, HumanMessage):
        st.chat_message("human", avatar="질문.png").write(msg.content)


# 사용자 입력 처리
if prompt := st.chat_input("부동산 시장에 관해 물어보세요"):
    # 사용자 질문 메시지 출력 및 저장
    with st.chat_message("human", avatar="질문.png"):
        formatted_prompt = prompt.replace("\n", "<br>")
        st.markdown(f'<div class="user-message">{formatted_prompt}</div>', unsafe_allow_html=True)
    st.session_state.file_qa_messages.append(HumanMessage(content=prompt))

    # 새로운 질문에 대한 응답 생성
    with st.chat_message("ai", avatar="챗봇.png"):
        response_container = st.empty()  # 여기서 새로운 빈 컨테이너 생성
        with st.spinner("답변생성중..."):
            response = get_chatbot_response(prompt, st.session_state.file_qa_content)

            if response:
                st.session_state.file_qa_messages.append(AIMessage(content=response))
                stream_response(response, response_container)  # 여기서 response_container를 인자로 전달
            else:
                response_container.error("Failed to get a response. Please try again.")

# # 피드백 수집
# if len(st.session_state.file_qa_messages) > 1:
#     last_message = st.session_state.file_qa_messages[-1]
#     if isinstance(last_message, AIMessage):
#         feedback = streamlit_feedback(
#             feedback_type="thumbs",
#             key=f"file_qa_feedback_{len(st.session_state.file_qa_messages)}"
#         )
#         if feedback:
#             score = 1 if feedback["score"] == "thumbsup" else 0
#             client = Client()
#             run_id = client.create_run(
#                 project_name=os.getenv("LANGCHAIN_PROJECT"),
#                 name="File Q&A Interaction",
#                 inputs={"messages": [msg.content for msg in st.session_state.file_qa_messages[1:]]},
#             ).id
#             client.create_feedback(run_id, "user_score", score=score)
#             st.toast("Feedback submitted!", icon="✅")