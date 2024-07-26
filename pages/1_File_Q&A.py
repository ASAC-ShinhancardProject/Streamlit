import streamlit as st
import anthropic
from langchain_anthropic import ChatAnthropic
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder, HumanMessagePromptTemplate
from langchain_core.messages import HumanMessage, AIMessage, SystemMessage
from langsmith import Client
from streamlit_feedback import streamlit_feedback
import os
from dotenv import load_dotenv
from PyPDF2 import PdfReader
import pandas as pd
from pptx import Presentation
import io

# 환경 변수 로드
load_dotenv()

# Streamlit 앱 설정
st.set_page_config(page_title="File Q&A and Chatbot", page_icon="📝", layout="wide")
st.title("📝 File Q&A and Chatbot")

# API 키 설정
anthropic_api_key = st.secrets.get("ANTHROPIC_API_KEY") or os.getenv("ANTHROPIC_API_KEY")

if not anthropic_api_key:
    st.error("Anthropic API key is not set. Please set it in your secrets or .env file.")
    st.stop()

# 세션 상태 초기화
if 'initialized' not in st.session_state:
    st.session_state.file_qa_messages = [SystemMessage(content="You are a helpful AI assistant.")]
    st.session_state.file_qa_content = None
    st.session_state.file_qa_data = None
    st.session_state.initialized = True

# 사이드바에 파일 업로드 위젯 추가
with st.sidebar:
    st.header("File Upload (Optional)")
    uploaded_file = st.file_uploader("Upload a file", type=("txt", "pdf", "csv", "xlsx", "pptx"))
    if uploaded_file:
        st.success(f"File '{uploaded_file.name}' uploaded successfully.")

# 파일 처리 함수
def process_file(file):
    if file is None:
        return None, None
    try:
        if file.type == "text/plain":
            content = file.read().decode()
            return content, None
        elif file.type == "application/pdf":
            pdf_reader = PdfReader(io.BytesIO(file.read()))
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text, None
        elif file.type == "text/csv":
            df = pd.read_csv(file)
            return df.to_string(), df
        elif file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            df = pd.read_excel(file)
            return df.to_string(), df
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = Presentation(file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + "\n"
            return text, None
        else:
            return f"Unsupported file type: {file.type}", None
    except Exception as e:
        st.error(f"Error processing file: {str(e)}")
        return None, None

# 파일 업로드 및 처리
if uploaded_file:
    with st.spinner("Processing uploaded file..."):
        st.session_state.file_qa_content, st.session_state.file_qa_data = process_file(uploaded_file)
        if st.session_state.file_qa_content:
            st.sidebar.success(f"File '{uploaded_file.name}' processed successfully.")
        else:
            st.sidebar.error("Failed to process the file.")

# 메인 영역을 세로로 배치
if st.session_state.file_qa_content:
    st.header("File Content")
    if isinstance(st.session_state.file_qa_data, pd.DataFrame):
        st.dataframe(st.session_state.file_qa_data)
    else:
        st.text_area("Content", st.session_state.file_qa_content[:1000] + "...", height=300)

st.header("Chat")

# 채팅 기록 표시
for message in st.session_state.file_qa_messages[1:]:
    with st.chat_message(message.type):
        st.markdown(message.content)

# 챗봇 응답 함수
def get_chatbot_response(prompt, file_content=None):
    try:
        client = anthropic.Anthropic(api_key=anthropic_api_key)
        if file_content:
            system_content = "You are an AI assistant that answers questions based on the given document and your general knowledge. If the information is not in the document, use your general knowledge to answer."
            user_content = f"Here's a document:\n\n<document>{file_content}</document>\n\nBased on this document and your general knowledge, answer the following question: {prompt}"
        else:
            system_content = "You are a helpful AI assistant."
            user_content = prompt
        
        response = client.messages.create(
            model="claude-3-sonnet-20240229",
            max_tokens=1000,
            system=system_content,
            messages=[
                {"role": "user", "content": user_content}
            ]
        )
        
        return response.content[0].text
    except Exception as e:
        st.error(f"Error in chatbot response: {str(e)}")
        return None

# 사용자 입력 처리
if prompt := st.chat_input("Ask a question or chat with the AI"):
    st.session_state.file_qa_messages.append(HumanMessage(content=prompt))
    with st.chat_message("human"):
        st.markdown(prompt)

    with st.chat_message("ai"):
        with st.spinner("Thinking..."):
            response = get_chatbot_response(prompt, st.session_state.file_qa_content)

            if response:
                st.markdown(response)
                st.session_state.file_qa_messages.append(AIMessage(content=response))
            else:
                st.error("Failed to get a response. Please try again.")

# 피드백 수집
if len(st.session_state.file_qa_messages) > 1:
    last_message = st.session_state.file_qa_messages[-1]
    if isinstance(last_message, AIMessage):
        feedback = streamlit_feedback(
            feedback_type="thumbs",
            key=f"file_qa_feedback_{len(st.session_state.file_qa_messages)}"
        )
        if feedback:
            score = 1 if feedback["score"] == "thumbsup" else 0
            client = Client()
            run_id = client.create_run(
                project_name=os.getenv("LANGCHAIN_PROJECT"),
                name="File Q&A Interaction",
                inputs={"messages": [msg.content for msg in st.session_state.file_qa_messages[1:]]},
            ).id
            client.create_feedback(run_id, "user_score", score=score)
            st.toast("Feedback submitted!", icon="✅")