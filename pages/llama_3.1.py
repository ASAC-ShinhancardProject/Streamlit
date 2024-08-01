import streamlit as st
import requests
import json
import subprocess
import io
from pptx import Presentation
from PyPDF2 import PdfReader

# 함수 정의
def read_ppt(file):
    prs = Presentation(file)
    text_content = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text_content.append(shape.text)
    return '\n'.join(text_content)

def read_pdf(file):
    pdf = PdfReader(file)
    text_content = []
    for page in pdf.pages:
        text_content.append(page.extract_text())
    return '\n'.join(text_content)

def check_ollama_status():
    try:
        result = subprocess.run(['ollama', 'list'], capture_output=True, text=True)
        if result.returncode == 0:
            return "Ollama가 정상적으로 실행 중입니다."
        else:
            return f"Ollama 실행 중 오류 발생: {result.stderr}"
    except FileNotFoundError:
        return "Ollama가 설치되어 있지 않거나 PATH에 없습니다."

def generate_llama_response(prompt_input):
    if "그려" in prompt_input or "그림" in prompt_input:
        return {
            "response": "죄송합니다. 저는 텍스트 기반 AI 어시스턴트로, 실제로 그림을 그리거나 이미지를 생성할 수 없습니다. 대신 그림에 대한 설명을 제공하거나 관련 정보를 알려드릴 수 있습니다. 어떤 도움이 필요하신가요?",
            "is_error": False
        }
    
    string_dialogue = f"You are a helpful assistant. Your responses should be in a {response_style.lower()} style. Respond in Korean. You are a text-based AI assistant and cannot generate, produce, edit, manipulate or create images.\n\n"
    for message in st.session_state.messages:
        if isinstance(message, dict):
            string_dialogue += f"{message['role'].capitalize()}: {message['content']}\n\n"
        elif hasattr(message, "type") and hasattr(message, "content"):
            string_dialogue += f"{message.type.capitalize()}: {message.content}\n\n"
    string_dialogue += f"Human: {prompt_input}\n\nAssistant: "
    
    data = {
        "model": selected_model,
        "prompt": string_dialogue,
        "stream": True,
        "temperature": temperature,
        "max_tokens": max_length,
    }
    
    try:
        response = requests.post(f"{OLLAMA_URL}/api/generate", json=data, stream=True, timeout=30)
        response.raise_for_status()
        return {"response": response, "is_error": False}
    except requests.exceptions.Timeout:
        return {"response": "API 요청 시간이 초과되었습니다. Ollama 서버의 상태를 확인해 주세요.", "is_error": True}
    except requests.exceptions.ConnectionError:
        return {"response": "Ollama 서버에 연결할 수 없습니다. 서버가 실행 중인지 확인해 주세요.", "is_error": True}
    except requests.exceptions.RequestException as e:
        return {"response": f"API 요청 중 오류가 발생했습니다: {str(e)}", "is_error": True}

def clear_chat_history():
    st.session_state.messages = [{"role": "assistant", "content": "무엇을 도와드릴까요?"}]
    if 'uploaded_file_content' in st.session_state:
        del st.session_state.uploaded_file_content
    if 'uploaded_file_name' in st.session_state:
        del st.session_state.uploaded_file_name

# App title
st.set_page_config(page_title="🦙💬 Llama Chatbot")

# Sidebar contents
with st.sidebar:
    st.title('🦙💬 Llama Chatbot')
    st.write('This chatbot uses the Llama model via Ollama.')
    
    st.subheader('Models and parameters')
    selected_model = st.selectbox('Choose a model', 
                                  ['llama3.1:latest', 'llama3:latest'],
                                  key='selected_model')
    temperature = st.slider('Temperature', min_value=0.1, max_value=2.0, value=0.1, step=0.1)
    max_length = st.slider('Max Length', min_value=64, max_value=4096, value=1024, step=64)
    response_style = st.selectbox('Response Style', ['Formal', 'Casual', 'Professional'], index=0)

# Ollama 서버 URL 설정
OLLAMA_URL = st.sidebar.text_input("Ollama 서버 URL", "http://localhost:11434")

# Store LLM generated responses
if "messages" not in st.session_state:
    st.session_state.messages = [{"role": "assistant", "content": "무엇을 도와드릴까요?"}]

# Display chat messages
for message in st.session_state.messages:
    if isinstance(message, dict):
        role = message.get("role", "assistant")
        content = message.get("content", "")
    elif hasattr(message, "type") and hasattr(message, "content"):
        role = message.type
        content = message.content
    else:
        st.warning(f"Unexpected message format: {message}")
        continue
    
    with st.chat_message(role):
        st.write(content)

st.sidebar.button('Clear Chat History', on_click=clear_chat_history, key='clear_chat_history')

# User-provided prompt
prompt = st.chat_input("메시지를 입력하세요")

# 파일 업로드 및 분석 기능
with st.expander("파일 업로드 및 분석", expanded=False):
    uploaded_file = st.file_uploader("파일을 업로드하세요", type=["txt", "pdf", "ppt", "pptx"])
    if uploaded_file is not None:
        file_contents = uploaded_file.read()
        
        if uploaded_file.type == "application/pdf":
            file_text = read_pdf(io.BytesIO(file_contents))
        elif uploaded_file.type in ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
            file_text = read_ppt(io.BytesIO(file_contents))
        else:
            # 다양한 인코딩 시도
            encodings = ['utf-8', 'euc-kr', 'cp949', 'iso-8859-1']
            file_text = None
            for encoding in encodings:
                try:
                    file_text = file_contents.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            
            # 모든 인코딩 시도 실패 시 바이너리 모드로 읽기
            if file_text is None:
                file_text = str(file_contents)
        
        st.session_state.uploaded_file_content = file_text
        st.session_state.uploaded_file_name = uploaded_file.name
        st.success(f"파일 '{uploaded_file.name}'이 성공적으로 업로드되었습니다.")
        
        # 파일 내용 미리보기
        with st.expander("파일 내용 미리보기", expanded=False):
            st.text_area("File Content", value=file_text, height=300, max_chars=None, key="file_content_preview")
        
    file_analysis = st.button("파일 분석")

if prompt or (file_analysis and 'uploaded_file_content' in st.session_state):
    if file_analysis and 'uploaded_file_content' in st.session_state:
        file_prompt = f"다음 파일 '{st.session_state.uploaded_file_name}'의 내용을 분석해주세요:\n\n{st.session_state.uploaded_file_content}\n\n사용자 질문: {prompt if prompt else '파일을 분석해주세요.'}"
        st.session_state.messages.append({"role": "user", "content": f"파일 '{st.session_state.uploaded_file_name}' 분석 요청\n사용자 질문: {prompt if prompt else '파일을 분석해주세요.'}"})
        with st.chat_message("user"):
            st.write(f"파일 '{st.session_state.uploaded_file_name}' 분석 요청")
            if prompt:
                st.write(f"사용자 질문: {prompt}")
        prompt = file_prompt
    elif prompt:
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.write(prompt)

    # Generate a new response
    with st.chat_message("assistant"):
        response_placeholder = st.empty()
        full_response = ""
        response_placeholder.markdown("답변 생성 중...")  # 답변 생성 중 표시
        result = generate_llama_response(prompt)
        if result["is_error"]:
            st.error(result["response"])
        else:
            response = result["response"]
            if isinstance(response, str):  # 직접 반환된 문자열 응답 처리
                full_response = response
                response_placeholder.markdown(full_response)
                st.session_state.messages.append({"role": "assistant", "content": full_response})
            else:  # API 응답 스트림 처리
                try:
                    for chunk in response.iter_lines():
                        if chunk:
                            chunk_data = json.loads(chunk.decode())
                            full_response += chunk_data['response']
                            response_placeholder.markdown(full_response)
                    
                    if full_response:
                        st.session_state.messages.append({"role": "assistant", "content": full_response})
                    else:
                        st.warning("모델이 응답을 생성하지 않았습니다. 다시 시도해 주세요.")
                except json.JSONDecodeError as e:
                    st.error(f"응답 데이터 파싱 중 오류가 발생했습니다: {str(e)}")
                except KeyError as e:
                    st.error(f"응답 데이터에 필요한 키가 없습니다: {str(e)}")
                except Exception as e:
                    st.error(f"응답 처리 중 오류 발생: {str(e)}")

# Ollama 상태 확인 및 표시
ollama_status = check_ollama_status()
st.sidebar.write(f"Ollama 상태: {ollama_status}")